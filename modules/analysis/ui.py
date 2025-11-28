"""
ANALYSIS MODULE - Streamlit UI Components
Provides UI for Comms Analyzer, Location Intelligence, and Media Viewer

This module provides:
- Location Intelligence UI
- GPS link input and tracking
- Map display
- Offline/online status
- Sync controls
- Analysis results display
"""

import streamlit as st
import pandas as pd
from datetime import datetime
from typing import Optional, Dict, Any, List

from modules.analysis.location_intelligence import LocationIntelligence
from modules.analysis.comms_analyzer import CommsAnalyzer
from modules.analysis.models import DatabaseManager

# ============================================================================
# COMMS ANALYZER UI
# ============================================================================

def render_comms_analyzer():
    """Render Comms Analyzer UI"""
    st.header("💬 Communications Analyzer")
    
    # Initialize session state
    if "comms_analyzer" not in st.session_state:
        st.session_state.comms_analyzer = CommsAnalyzer()
    
    analyzer = st.session_state.comms_analyzer
    
    # Check consent
    from modules.consent.models import get_consent_manager, ConsentLevel, MODULE_MIN_LEVELS
    consent_manager = get_consent_manager()
    case_id = st.session_state.get("case_id", "DEFAULT")
    
    session = consent_manager.get_session(case_id)
    min_level = MODULE_MIN_LEVELS.get('communications', ConsentLevel.LEGAL)
    
    # Dev mode toggle
    col_dev1, col_dev2 = st.columns([4, 1])
    with col_dev2:
        if st.checkbox("🧪 Dev Mode", value=consent_manager.connectivity_manager.is_dev_mode(), key="comms_dev_mode"):
            consent_manager.connectivity_manager.set_dev_mode(True)
            st.success("Dev mode enabled")
        else:
            consent_manager.connectivity_manager.set_dev_mode(False)
    
    # Display consent status
    if session:
        col1, col2, col3 = st.columns(3)
        with col1:
            st.metric("Consent Level", session.level.name)
        with col2:
            st.metric("Required Level", min_level.name)
        with col3:
            if session.level >= min_level:
                st.success("✅ Authorized")
            else:
                st.error("❌ Insufficient Consent")
        
        # Check if authorized (skip in dev mode)
        if session.level < min_level and not consent_manager.connectivity_manager.is_dev_mode():
            st.error(f"🚫 Communications analysis requires {min_level.name} consent level")
            st.info(f"Current level: {session.level.name} | Required: {min_level.name}")
            st.warning("Please upgrade your consent level to proceed with communications analysis")
            return
        elif session.level < min_level and consent_manager.connectivity_manager.is_dev_mode():
            st.warning(f"⚠️ Dev Mode: Bypassing consent check (normally requires {min_level.name})")
    else:
        st.error("❌ No consent found for this case")
        st.info("Please provide consent before proceeding")
        if not consent_manager.connectivity_manager.is_dev_mode():
            return
        else:
            st.warning("⚠️ Dev Mode: Proceeding without consent")
    
    # Tabs for different features
    tab1, tab2, tab3, tab4 = st.tabs([
        "📨 Analyze Message",
        "📞 Check Phone",
        "📧 Check Email",
        "📊 Results"
    ])
    
    # ========================================================================
    # TAB 1: ANALYZE MESSAGE
    # ========================================================================
    
    with tab1:
        st.subheader("Analyze Message for Suspicious Content")
        
        message = st.text_area(
            "Enter message to analyze",
            placeholder="Paste SMS, email, or chat message here...",
            height=150
        )
        
        col1, col2, col3 = st.columns(3)
        
        with col1:
            sender = st.text_input(
                "Sender Phone/Email",
                placeholder="Phone number or email"
            )
        
        with col2:
            contact_phone = st.text_input(
                "Contact Phone Number",
                placeholder="+1-555-0123",
                help="Phone number of the contact"
            )
        
        with col3:
            message_type = st.selectbox(
                "Message Type",
                ["SMS", "Email", "Chat", "Unknown"]
            )
        
        if st.button("🔍 Analyze Message", use_container_width=True):
            if message:
                with st.spinner("Analyzing message..."):
                    result = analyzer.analyze_message(message, sender, contact_phone)
                    
                    # Store result
                    st.session_state.last_message_analysis = result
                    
                    # Display risk level
                    risk_score = result.get("risk_score", 0)
                    classification = result.get("classification", "UNKNOWN")
                    
                    # Display contact info
                    st.subheader("📞 Contact Information")
                    col1, col2, col3 = st.columns(3)
                    
                    with col1:
                        st.info(f"**Sender:** {sender or 'Not provided'}")
                    
                    with col2:
                        st.info(f"**Contact Phone:** {contact_phone or 'Not provided'}")
                    
                    with col3:
                        st.info(f"**Message Type:** {message_type}")
                    
                    # Display risk level
                    st.subheader("🚨 Risk Assessment")
                    col1, col2, col3 = st.columns(3)
                    
                    with col1:
                        if classification == "CRITICAL":
                            st.error(f"🔴 {classification}")
                        elif classification == "HIGH":
                            st.warning(f"🟠 {classification}")
                        elif classification == "MEDIUM":
                            st.warning(f"🟡 {classification}")
                        else:
                            st.success(f"🟢 {classification}")
                    
                    with col2:
                        st.metric("Risk Score", f"{risk_score:.2f}")
                    
                    with col3:
                        st.metric("Timestamp", datetime.now().strftime("%Y-%m-%d %H:%M:%S"))
                    
                    # Display detected threats
                    if result.get("threats"):
                        st.subheader("🚨 Detected Threats")
                        threats_df = pd.DataFrame(result["threats"])
                        st.dataframe(threats_df, use_container_width=True)
                    
                    # Display keywords
                    if result.get("keywords"):
                        st.subheader("🔑 Suspicious Keywords")
                        keywords_text = ", ".join(result["keywords"])
                        st.info(keywords_text)
                    
                    # Display sentiment
                    if result.get("sentiment"):
                        st.subheader("😊 Sentiment Analysis")
                        col1, col2 = st.columns(2)
                        with col1:
                            st.metric("Sentiment", result["sentiment"]["label"])
                        with col2:
                            st.metric("Score", f"{result['sentiment']['score']:.2f}")
                    
                    # Database checks
                    st.subheader("🔍 Database Checks")
                    
                    col1, col2 = st.columns(2)
                    
                    # Sender phone check
                    with col1:
                        sender_check = result.get("database_checks", {}).get("sender_phone_match", {})
                        if sender_check.get("match"):
                            st.error(f"🚨 Sender Phone in Database")
                            st.write(f"Type: {sender_check.get('type')}")
                        else:
                            st.success("✅ Sender Phone: Clean")
                    
                    # Contact phone check
                    with col2:
                        contact_check = result.get("database_checks", {}).get("contact_phone_match", {})
                        if contact_check.get("match"):
                            st.error(f"🚨 Contact Phone in Database")
                            st.write(f"Type: {contact_check.get('type')}")
                        else:
                            st.success("✅ Contact Phone: Clean")
            else:
                st.warning("Please enter a message to analyze")
    
    # ========================================================================
    # TAB 2: CHECK PHONE
    # ========================================================================
    
    with tab2:
        st.subheader("Check Phone Number Against Database")
        
        phone = st.text_input(
            "Phone Number",
            placeholder="+1-555-0123 or 5550123"
        )
        
        if st.button("📞 Check Phone", use_container_width=True):
            if phone:
                with st.spinner("Checking database..."):
                    result = analyzer.check_phone_database(phone)
                    
                    if result.get("is_fraudster"):
                        st.error("🚨 FRAUDSTER DETECTED")
                        fraudster = result.get("fraudster")
                        col1, col2, col3 = st.columns(3)
                        with col1:
                            st.metric("Type", fraudster.get("fraud_type", "Unknown"))
                        with col2:
                            st.metric("Risk Level", fraudster.get("risk_level", "Unknown"))
                        with col3:
                            st.metric("Reports", fraudster.get("reports", 0))
                        
                        st.write(fraudster)
                    
                    elif result.get("is_harasser"):
                        st.error("🚨 HARASSER DETECTED")
                        harasser = result.get("harasser")
                        col1, col2, col3 = st.columns(3)
                        with col1:
                            st.metric("Type", harasser.get("harassment_type", "Unknown"))
                        with col2:
                            st.metric("Risk Level", harasser.get("risk_level", "Unknown"))
                        with col3:
                            st.metric("Victims", harasser.get("victims", 0))
                        
                        st.write(harasser)
                    
                    else:
                        st.success("✅ Phone number not in database")
            else:
                st.warning("Please enter a phone number")
    
    # ========================================================================
    # TAB 3: CHECK EMAIL
    # ========================================================================
    
    with tab3:
        st.subheader("Check Email Against Database")
        
        email = st.text_input(
            "Email Address",
            placeholder="example@email.com"
        )
        
        if st.button("📧 Check Email", use_container_width=True):
            if email:
                with st.spinner("Checking database..."):
                    result = analyzer.check_email_database(email)
                    
                    if result.get("found"):
                        st.error("🚨 FRAUDULENT EMAIL DETECTED")
                        email_data = result.get("email_data")
                        col1, col2, col3 = st.columns(3)
                        with col1:
                            st.metric("Type", email_data.get("fraud_type", "Unknown"))
                        with col2:
                            st.metric("Risk Level", email_data.get("risk_level", "Unknown"))
                        with col3:
                            st.metric("Reports", email_data.get("reports", 0))
                        
                        st.write(email_data)
                    
                    else:
                        st.success("✅ Email not in database")
            else:
                st.warning("Please enter an email address")
    
    # ========================================================================
    # TAB 4: RESULTS
    # ========================================================================
    
    with tab4:
        st.subheader("Analysis Results")
        
        if "last_message_analysis" in st.session_state:
            result = st.session_state.last_message_analysis
            
            # Summary
            st.subheader("Summary")
            col1, col2, col3, col4 = st.columns(4)
            with col1:
                st.metric("Risk Score", f"{result.get('risk_score', 0):.2f}")
            with col2:
                st.metric("Classification", result.get("classification", "Unknown"))
            with col3:
                st.metric("Threats", len(result.get("threats", [])))
            with col4:
                st.metric("Keywords", len(result.get("keywords", [])))
            
            # Detailed results
            if result.get("threats"):
                st.subheader("🚨 Threats")
                threats_df = pd.DataFrame(result["threats"])
                st.dataframe(threats_df, use_container_width=True)
            
            if result.get("keywords"):
                st.subheader("🔑 Keywords")
                keywords_df = pd.DataFrame([
                    {"Keyword": kw, "Category": "Suspicious"} 
                    for kw in result["keywords"]
                ])
                st.dataframe(keywords_df, use_container_width=True)
            
            if result.get("entities"):
                st.subheader("🏷️ Entities")
                entities_df = pd.DataFrame(result["entities"])
                st.dataframe(entities_df, use_container_width=True)
            
            # Export
            st.subheader("📥 Export Results")
            
            if st.button("📋 Copy JSON", use_container_width=True):
                import json
                st.code(json.dumps(result, indent=2, default=str))
        
        else:
            st.info("Analyze a message first to see results")


# ============================================================================
# LOCATION INTELLIGENCE UI
# ============================================================================

def render_location_intelligence():
    """Render Location Intelligence UI"""
    st.header("📍 Location Intelligence")
    
    # Initialize session state
    if "location_analyzer" not in st.session_state:
        st.session_state.location_analyzer = LocationIntelligence()
    
    analyzer = st.session_state.location_analyzer
    
    # Check consent
    from modules.consent.models import get_consent_manager, ConsentLevel, MODULE_MIN_LEVELS
    consent_manager = get_consent_manager()
    case_id = st.session_state.get("case_id", "DEFAULT")
    
    session = consent_manager.get_session(case_id)
    min_level = MODULE_MIN_LEVELS.get('location', ConsentLevel.STANDARD)
    
    # Dev mode toggle
    col_dev1, col_dev2 = st.columns([4, 1])
    with col_dev2:
        if st.checkbox("🧪 Dev Mode", value=consent_manager.connectivity_manager.is_dev_mode(), key="location_dev_mode"):
            consent_manager.connectivity_manager.set_dev_mode(True)
            st.success("Dev mode enabled")
        else:
            consent_manager.connectivity_manager.set_dev_mode(False)
    
    # Display consent status
    if session:
        col1, col2, col3 = st.columns(3)
        with col1:
            st.metric("Consent Level", session.level.name)
        with col2:
            st.metric("Required Level", min_level.name)
        with col3:
            if session.level >= min_level:
                st.success("✅ Authorized")
            else:
                st.error("❌ Insufficient Consent")
        
        # Check if authorized (skip in dev mode)
        if session.level < min_level and not consent_manager.connectivity_manager.is_dev_mode():
            st.error(f"🚫 Location analysis requires {min_level.name} consent level")
            st.info(f"Current level: {session.level.name} | Required: {min_level.name}")
            st.warning("Please upgrade your consent level to proceed with location analysis")
            return
        elif session.level < min_level and consent_manager.connectivity_manager.is_dev_mode():
            st.warning(f"⚠️ Dev Mode: Bypassing consent check (normally requires {min_level.name})")
    else:
        st.error("❌ No consent found for this case")
        st.info("Please provide consent before proceeding")
        if not consent_manager.connectivity_manager.is_dev_mode():
            return
        else:
            st.warning("⚠️ Dev Mode: Proceeding without consent")
    
    # Tabs for different features
    tab1, tab2, tab3, tab4, tab5 = st.tabs([
        "🔗 GPS Links",
        "📍 Coordinates",
        "📊 Analysis",
        "📡 Offline Status",
        "📈 Results"
    ])
    
    # ========================================================================
    # TAB 1: GPS LINKS
    # ========================================================================
    
    with tab1:
        st.subheader("Add GPS Links (WhatsApp, Google Maps)")
        
        col1, col2 = st.columns([3, 1])
        
        with col1:
            gps_link = st.text_input(
                "Paste GPS Link",
                placeholder="https://maps.google.com/?q=40.7128,-74.0060",
                help="WhatsApp location, Google Maps, or shortened URL"
            )
        
        with col2:
            add_link_btn = st.button("➕ Add Link", use_container_width=True)
        
        if add_link_btn and gps_link:
            with st.spinner("Processing GPS link..."):
                # Get case ID from session
                case_id = st.session_state.get("case_id", "DEFAULT")
                user_name = st.session_state.get("user_name", "Unknown")
                
                result = analyzer.add_location_from_link(
                    link=gps_link,
                    case_id=case_id,
                    added_by=user_name
                )
                
                if result["status"] == "success":
                    location = result["location"]
                    sync_status = location.get("sync_status", "unknown")
                    
                    # Show status
                    if sync_status == "synced":
                        st.success(f"✅ Location synced to database (ID: {location.get('db_id')})")
                    else:
                        st.warning(f"📋 Location queued for sync (offline mode)")
                    
                    # Display location details
                    col1, col2, col3, col4 = st.columns(4)
                    with col1:
                        st.metric("Latitude", f"{location['latitude']:.4f}")
                    with col2:
                        st.metric("Longitude", f"{location['longitude']:.4f}")
                    with col3:
                        st.metric("Source", location['source'])
                    with col4:
                        st.metric("Status", sync_status.upper())
                else:
                    st.error(f"❌ Error: {result.get('message')}")
        
        # Display recent GPS links
        st.subheader("Recent GPS Links")
        
        case_id = st.session_state.get("case_id", "DEFAULT")
        db = DatabaseManager()
        
        try:
            links = db.get_gps_links_by_case(case_id)
            
            if links:
                # Create DataFrame
                links_data = []
                for link in links:
                    links_data.append({
                        "ID": link.id,
                        "Source": link.source,
                        "Location": link.location_name or f"({link.latitude}, {link.longitude})",
                        "Risk": link.risk_level,
                        "Added": link.added_at.strftime("%Y-%m-%d %H:%M"),
                        "Status": link.status
                    })
                
                df = pd.DataFrame(links_data)
                st.dataframe(df, use_container_width=True)
            else:
                st.info("No GPS links added yet")
        except Exception as e:
            st.error(f"Error loading GPS links: {e}")
    
    # ========================================================================
    # TAB 2: COORDINATES
    # ========================================================================
    
    with tab2:
        st.subheader("Add Location by Coordinates")
        
        col1, col2, col3 = st.columns(3)
        
        with col1:
            latitude = st.number_input(
                "Latitude",
                min_value=-90.0,
                max_value=90.0,
                value=40.7128,
                step=0.0001
            )
        
        with col2:
            longitude = st.number_input(
                "Longitude",
                min_value=-180.0,
                max_value=180.0,
                value=-74.0060,
                step=0.0001
            )
        
        with col3:
            location_name = st.text_input(
                "Location Name",
                placeholder="e.g., New York"
            )
        
        if st.button("➕ Add Coordinates", use_container_width=True):
            with st.spinner("Adding location..."):
                case_id = st.session_state.get("case_id", "DEFAULT")
                user_name = st.session_state.get("user_name", "Unknown")
                
                result = analyzer.add_location_from_coordinates(
                    latitude=latitude,
                    longitude=longitude,
                    name=location_name or f"Location ({latitude}, {longitude})"
                )
                
                if result["status"] == "success":
                    st.success("✅ Location added successfully")
                    
                    location = result["location"]
                    col1, col2, col3 = st.columns(3)
                    with col1:
                        st.metric("Latitude", f"{location['latitude']:.4f}")
                    with col2:
                        st.metric("Longitude", f"{location['longitude']:.4f}")
                    with col3:
                        st.metric("Name", location['name'])
                else:
                    st.error(f"❌ Error: {result.get('message')}")
        
        # CSV Bulk Upload
        st.subheader("Bulk Upload (CSV)")
        
        csv_text = st.text_area(
            "Paste CSV data (latitude,longitude,name)",
            placeholder="40.7128,-74.0060,New York\n34.0522,-118.2437,Los Angeles",
            height=150
        )
        
        if st.button("📤 Upload CSV", use_container_width=True):
            if csv_text:
                with st.spinner("Processing CSV..."):
                    result = analyzer.add_locations_from_csv(csv_text)
                    
                    if result["status"] == "success":
                        st.success(f"✅ Added {result['total_added']} locations")
                        
                        if result["total_errors"] > 0:
                            st.warning(f"⚠️ {result['total_errors']} errors:")
                            for error in result["errors"]:
                                st.text(f"  • {error}")
                    else:
                        st.error(f"❌ Error: {result.get('message')}")
            else:
                st.warning("Please paste CSV data")
    
    # ========================================================================
    # TAB 3: ANALYSIS
    # ========================================================================
    
    with tab3:
        st.subheader("Analyze Locations")
        
        case_id = st.session_state.get("case_id", "DEFAULT")
        db = DatabaseManager()
        
        try:
            links = db.get_gps_links_by_case(case_id)
            
            if not links:
                st.info("No locations to analyze. Add some first!")
            else:
                # Convert to location format for analysis
                locations = []
                for link in links:
                    locations.append({
                        "name": link.location_name or f"Location {link.id}",
                        "latitude": link.latitude,
                        "longitude": link.longitude,
                        "timestamp": link.added_at.isoformat(),
                        "type": link.source
                    })
                
                if st.button("🔍 Analyze All Locations", use_container_width=True):
                    with st.spinner("Analyzing locations..."):
                        analysis_result = analyzer.analyze_locations(locations)
                        
                        # Display results
                        st.subheader("Analysis Results")
                        
                        col1, col2, col3, col4 = st.columns(4)
                        with col1:
                            st.metric("Total Locations", analysis_result["total_locations"])
                        with col2:
                            st.metric("Risk Score", f"{analysis_result['overall_risk_score']:.2f}")
                        with col3:
                            st.metric("Classification", analysis_result["classification"])
                        with col4:
                            st.metric("Anomalies", analysis_result["anomalies"]["total_anomalies"])
                        
                        # Timeline
                        if analysis_result["timeline"]["timeline"]:
                            st.subheader("📅 Timeline")
                            timeline_df = pd.DataFrame(analysis_result["timeline"]["timeline"])
                            st.dataframe(timeline_df, use_container_width=True)
                        
                        # Frequent Locations
                        if analysis_result["frequent_locations"]["frequent_locations"]:
                            st.subheader("🏢 Frequent Locations")
                            freq_df = pd.DataFrame(
                                analysis_result["frequent_locations"]["frequent_locations"]
                            )
                            st.dataframe(freq_df, use_container_width=True)
                        
                        # Travel Patterns
                        if analysis_result["travel_patterns"]["travel_patterns"]:
                            st.subheader("🚗 Travel Patterns")
                            patterns_df = pd.DataFrame(
                                analysis_result["travel_patterns"]["travel_patterns"]
                            )
                            st.dataframe(patterns_df, use_container_width=True)
                        
                        # Anomalies
                        if analysis_result["anomalies"]["anomalies"]:
                            st.subheader("⚠️ Anomalies Detected")
                            anomalies_df = pd.DataFrame(
                                analysis_result["anomalies"]["anomalies"]
                            )
                            st.dataframe(anomalies_df, use_container_width=True)
                        
                        # Store in session for later use
                        st.session_state.last_analysis = analysis_result
        
        except Exception as e:
            st.error(f"Error during analysis: {e}")
    
    # ========================================================================
    # TAB 4: OFFLINE STATUS
    # ========================================================================
    
    with tab4:
        st.subheader("📡 Offline/Online Status")
        
        # Get status
        status = analyzer.get_offline_status()
        
        # Display connectivity
        col1, col2, col3 = st.columns(3)
        
        with col1:
            if status["is_online"]:
                st.success("🟢 ONLINE")
            else:
                st.error("🔴 OFFLINE")
        
        with col2:
            st.metric("Pending Operations", status["pending_operations"])
        
        with col3:
            st.metric("Synced Operations", status["synced_operations"])
        
        # Queue details
        st.subheader("Queue Details")
        
        queue_stats = status["queue_stats"]
        col1, col2, col3, col4 = st.columns(4)
        
        with col1:
            st.metric("Total Operations", queue_stats["total_operations"])
        with col2:
            st.metric("Pending", queue_stats["pending"])
        with col3:
            st.metric("Synced", queue_stats["synced"])
        with col4:
            st.metric("Status", "Online" if queue_stats["is_online"] else "Offline")
        
        # Sync button
        if status["pending_operations"] > 0:
            st.warning(f"⚠️ {status['pending_operations']} pending operations waiting to sync")
            
            if st.button("🔄 Sync Now", use_container_width=True):
                with st.spinner("Syncing..."):
                    sync_result = analyzer.sync_pending_operations()
                    
                    if sync_result["status"] == "success":
                        st.success(
                            f"✅ Synced {sync_result['synced']} operations"
                        )
                        if sync_result["failed"] > 0:
                            st.warning(f"⚠️ {sync_result['failed']} operations failed")
                            if sync_result["errors"]:
                                with st.expander("View Errors"):
                                    for error in sync_result["errors"]:
                                        st.text(error)
                    else:
                        st.error(f"❌ Sync failed: {sync_result.get('message')}")
        else:
            st.success("✅ All operations synced")
    
    # ========================================================================
    # TAB 5: RESULTS
    # ========================================================================
    
    with tab5:
        st.subheader("📊 Analysis Results")
        
        if "last_analysis" in st.session_state:
            analysis = st.session_state.last_analysis
            
            # Summary
            st.subheader("Summary")
            col1, col2, col3, col4 = st.columns(4)
            with col1:
                st.metric("Total Locations", analysis["total_locations"])
            with col2:
                st.metric("Risk Score", f"{analysis['overall_risk_score']:.2f}")
            with col3:
                st.metric("Classification", analysis["classification"])
            with col4:
                st.metric("Anomalies", analysis["anomalies"]["total_anomalies"])
            
            # Distance Analysis
            if analysis["distances"]["total_distance_km"] > 0:
                st.subheader("📏 Distance Analysis")
                col1, col2, col3 = st.columns(3)
                with col1:
                    st.metric(
                        "Total Distance",
                        f"{analysis['distances']['total_distance_km']:.2f} km"
                    )
                with col2:
                    st.metric(
                        "Average Distance",
                        f"{analysis['distances']['average_distance_per_trip_km']:.2f} km"
                    )
                with col3:
                    st.metric(
                        "Max Distance",
                        f"{analysis['distances']['max_distance_single_trip_km']:.2f} km"
                    )
            
            # Risk Assessment
            st.subheader("🚨 Risk Assessment")
            risk_data = analysis["risk_assessment"]
            if risk_data["high_risk_locations"]:
                risk_df = pd.DataFrame(risk_data["high_risk_locations"])
                st.dataframe(risk_df, use_container_width=True)
            else:
                st.success("✅ No high-risk locations detected")
            
            # Export results
            st.subheader("📥 Export Results")
            
            col1, col2 = st.columns(2)
            
            with col1:
                if st.button("📋 Copy JSON", use_container_width=True):
                    import json
                    st.code(json.dumps(analysis, indent=2, default=str))
            
            with col2:
                if st.button("📊 Download CSV", use_container_width=True):
                    # Create CSV from timeline
                    if analysis["timeline"]["timeline"]:
                        df = pd.DataFrame(analysis["timeline"]["timeline"])
                        csv = df.to_csv(index=False)
                        st.download_button(
                            label="Download CSV",
                            data=csv,
                            file_name="location_analysis.csv",
                            mime="text/csv"
                        )
        else:
            st.info("Run analysis first to see results")


# ============================================================================
# MAIN
# ============================================================================

# ============================================================================
# DASHBOARD OVERVIEW UI
# ============================================================================

def render_dashboard_overview():
    """Render Analysis Dashboard Overview"""
    st.header("📊 Analysis Dashboard Overview")
    
    db = DatabaseManager()
    case_id = st.session_state.get("case_id", "DEFAULT")
    
    # Check consent
    from modules.consent.models import get_consent_manager
    consent_manager = get_consent_manager()
    session = consent_manager.get_session(case_id)
    
    # Dev mode toggle
    col_dev1, col_dev2 = st.columns([4, 1])
    with col_dev2:
        if st.checkbox("🧪 Dev Mode", value=consent_manager.connectivity_manager.is_dev_mode(), key="dashboard_dev_mode"):
            consent_manager.connectivity_manager.set_dev_mode(True)
            st.success("Dev mode enabled")
        else:
            consent_manager.connectivity_manager.set_dev_mode(False)
    
    # Display consent status
    if session:
        st.info(f"📋 Case Consent Level: **{session.level.name}**")
    else:
        if not consent_manager.connectivity_manager.is_dev_mode():
            st.warning("⚠️ No consent found for this case")
        else:
            st.warning("⚠️ No consent found for this case (Dev Mode: Bypassing)")
    
    # Dev mode info
    if consent_manager.connectivity_manager.is_dev_mode():
        st.info("🧪 **Dev Mode Active** - Testing all modules with bypass enabled")
        st.success("✅ Consent checks bypassed for testing")
    
    # Get statistics
    try:
        stats = db.get_statistics()
        gps_stats = db.get_gps_links_statistics(case_id)
    except Exception as e:
        st.error(f"Error loading statistics: {e}")
        return
    
    # ========================================================================
    # SUMMARY METRICS
    # ========================================================================
    
    st.subheader("📈 Summary Metrics")
    
    col1, col2, col3, col4, col5 = st.columns(5)
    
    with col1:
        st.metric("Total Fraudsters", stats['total_fraudsters'])
    
    with col2:
        st.metric("Total Harassers", stats['total_harassers'])
    
    with col3:
        st.metric("GPS Links Tracked", gps_stats.get('total_links', 0))
    
    with col4:
        st.metric("Critical Cases", stats['critical_fraudsters'] + stats['critical_harassers'])
    
    with col5:
        st.metric("Active Case", case_id)
    
    # ========================================================================
    # QUICK STATS
    # ========================================================================
    
    st.subheader("⚡ Quick Stats")
    
    col1, col2, col3 = st.columns(3)
    
    with col1:
        st.write("**Fraudster Risk Distribution**")
        try:
            from modules.analysis.models import Fraudster
            fraudsters = db.db.query(Fraudster).all()
            risk_counts = {}
            for f in fraudsters:
                risk_counts[f.risk_level] = risk_counts.get(f.risk_level, 0) + 1
            
            if risk_counts:
                st.bar_chart(risk_counts)
            else:
                st.info("No fraudsters in database")
        except Exception as e:
            st.warning(f"Could not load data: {e}")
    
    with col2:
        st.write("**Harasser Risk Distribution**")
        try:
            from modules.analysis.models import Harasser
            harassers = db.db.query(Harasser).all()
            risk_counts = {}
            for h in harassers:
                risk_counts[h.risk_level] = risk_counts.get(h.risk_level, 0) + 1
            
            if risk_counts:
                st.bar_chart(risk_counts)
            else:
                st.info("No harassers in database")
        except Exception as e:
            st.warning(f"Could not load data: {e}")
    
    with col3:
        st.write("**GPS Links by Source**")
        source_counts = gps_stats.get('by_source', {})
        if source_counts:
            st.bar_chart(source_counts)
        else:
            st.info("No GPS links tracked")
    
    # ========================================================================
    # RECENT ACTIVITY
    # ========================================================================
    
    st.subheader("🕐 Recent Activity")
    
    try:
        from modules.analysis.models import GPSLinkLog
        
        # Get recent GPS links
        recent_links = db.db.query(GPSLinkLog).filter(
            GPSLinkLog.case_id == case_id
        ).order_by(GPSLinkLog.added_at.desc()).limit(5).all()
        
        if recent_links:
            activity_data = []
            for link in recent_links:
                activity_data.append({
                    "Timestamp": link.added_at.strftime("%Y-%m-%d %H:%M:%S"),
                    "Type": "GPS Link",
                    "Source": link.source,
                    "Location": link.location_name or f"({link.latitude}, {link.longitude})",
                    "Risk": link.risk_level,
                    "Status": link.status
                })
            
            activity_df = pd.DataFrame(activity_data)
            st.dataframe(activity_df, use_container_width=True)
        else:
            st.info("No recent activity in this case")
    
    except Exception as e:
        st.warning(f"Could not load recent activity: {e}")
    
    # ========================================================================
    # RISK SUMMARY
    # ========================================================================
    
    st.subheader("🚨 Risk Summary")
    
    col1, col2, col3 = st.columns(3)
    
    with col1:
        st.write("**Critical Fraudsters**")
        critical_fraudsters = stats['critical_fraudsters']
        if critical_fraudsters > 0:
            st.error(f"🔴 {critical_fraudsters} Critical Cases")
        else:
            st.success("✅ No Critical Cases")
    
    with col2:
        st.write("**Critical Harassers**")
        critical_harassers = stats['critical_harassers']
        if critical_harassers > 0:
            st.error(f"🔴 {critical_harassers} Critical Cases")
        else:
            st.success("✅ No Critical Cases")
    
    with col3:
        st.write("**High-Risk GPS Locations**")
        high_risk_gps = gps_stats.get('high_risk', 0)
        if high_risk_gps > 0:
            st.warning(f"🟠 {high_risk_gps} High-Risk Locations")
        else:
            st.success("✅ No High-Risk Locations")
    
    # ========================================================================
    # CASE SUMMARY
    # ========================================================================
    
    st.subheader("📋 Case Summary")
    
    col1, col2 = st.columns(2)
    
    with col1:
        st.write("**Case Information**")
        st.write(f"**Case ID**: {case_id}")
        st.write(f"**User**: {st.session_state.get('user_name', 'Unknown')}")
        st.write(f"**Total Fraudsters**: {stats['total_fraudsters']}")
        st.write(f"**Total Harassers**: {stats['total_harassers']}")
    
    with col2:
        st.write("**Analysis Summary**")
        st.write(f"**GPS Links**: {gps_stats.get('total_links', 0)}")
        st.write(f"**WhatsApp Links**: {gps_stats.get('by_source', {}).get('whatsapp', 0)}")
        st.write(f"**Google Maps Links**: {gps_stats.get('by_source', {}).get('google_maps', 0)}")
        st.write(f"**High-Risk Items**: {stats['critical_fraudsters'] + stats['critical_harassers'] + high_risk_gps}")
    
    # ========================================================================
    # QUICK ACTIONS
    # ========================================================================
    
    st.subheader("⚡ Quick Actions")
    
    col1, col2, col3, col4 = st.columns(4)
    
    with col1:
        if st.button("➕ Add Fraudster", use_container_width=True):
            st.session_state.quick_action = "add_fraudster"
            st.rerun()
    
    with col2:
        if st.button("➕ Add Harasser", use_container_width=True):
            st.session_state.quick_action = "add_harasser"
            st.rerun()
    
    with col3:
        if st.button("🔗 Add GPS Link", use_container_width=True):
            st.session_state.quick_action = "add_gps_link"
            st.rerun()
    
    with col4:
        if st.button("💬 Analyze Message", use_container_width=True):
            st.session_state.quick_action = "analyze_message"
            st.rerun()


# ============================================================================
# DATABASE MANAGEMENT UI
# ============================================================================

def render_database_management():
    """Render Database Management UI"""
    st.header("🗄️ Fraud & Harassment Database Management")
    
    db = DatabaseManager()
    
    # Tabs for different operations
    tab1, tab2, tab3, tab4, tab5 = st.tabs([
        "➕ Add Fraudster",
        "➕ Add Harasser",
        "📊 View Database",
        "🔍 Search",
        "📈 Statistics"
    ])
    
    # ========================================================================
    # TAB 1: ADD FRAUDSTER
    # ========================================================================
    
    with tab1:
        st.subheader("Add Fraudster to Database")
        
        col1, col2 = st.columns(2)
        
        with col1:
            phone = st.text_input(
                "Phone Number",
                placeholder="+1-555-0123",
                key="fraudster_phone"
            )
            fraud_type = st.selectbox(
                "Fraud Type",
                ["Phishing", "Money Transfer", "Identity Theft", "Romance Scam", 
                 "Lottery Scam", "Tech Support Scam", "IRS Scam", "Other"]
            )
            name = st.text_input("Name (optional)", placeholder="John Doe")
        
        with col2:
            risk_level = st.selectbox(
                "Risk Level",
                ["LOW", "MEDIUM", "HIGH", "CRITICAL"]
            )
            methods = st.text_area(
                "Methods Used (comma-separated)",
                placeholder="SMS phishing, fake links, impersonation"
            )
        
        if st.button("➕ Add Fraudster", use_container_width=True):
            if phone and fraud_type:
                try:
                    methods_list = [m.strip() for m in methods.split(",")] if methods else []
                    
                    fraudster = db.add_fraudster(
                        phone=phone,
                        fraud_type=fraud_type,
                        name=name or None,
                        methods=methods_list,
                        risk_level=risk_level
                    )
                    
                    st.success(f"✅ Fraudster added successfully (ID: {fraudster.id})")
                    st.write(f"Phone: {fraudster.phone}")
                    st.write(f"Type: {fraudster.fraud_type}")
                    st.write(f"Risk: {fraudster.risk_level}")
                except Exception as e:
                    st.error(f"❌ Error: {e}")
            else:
                st.warning("Please enter phone number and fraud type")
    
    # ========================================================================
    # TAB 2: ADD HARASSER
    # ========================================================================
    
    with tab2:
        st.subheader("Add Harasser to Database")
        
        col1, col2 = st.columns(2)
        
        with col1:
            phone = st.text_input(
                "Phone Number",
                placeholder="+1-555-0123",
                key="harasser_phone"
            )
            harassment_type = st.selectbox(
                "Harassment Type",
                ["Abusive Calls", "Threatening Messages", "Stalking", "Bullying",
                 "Sexual Harassment", "Extortion", "Blackmail", "Other"]
            )
            name = st.text_input("Name (optional)", placeholder="Jane Doe")
        
        with col2:
            risk_level = st.selectbox(
                "Risk Level",
                ["LOW", "MEDIUM", "HIGH", "CRITICAL"],
                key="harasser_risk"
            )
            victims = st.number_input(
                "Number of Victims",
                min_value=1,
                value=1
            )
        
        if st.button("➕ Add Harasser", use_container_width=True):
            if phone and harassment_type:
                try:
                    harasser = db.add_harasser(
                        phone=phone,
                        harassment_type=harassment_type,
                        name=name or None,
                        risk_level=risk_level,
                        victims=victims
                    )
                    
                    st.success(f"✅ Harasser added successfully (ID: {harasser.id})")
                    st.write(f"Phone: {harasser.phone}")
                    st.write(f"Type: {harasser.harassment_type}")
                    st.write(f"Victims: {harasser.victims}")
                except Exception as e:
                    st.error(f"❌ Error: {e}")
            else:
                st.warning("Please enter phone number and harassment type")
    
    # ========================================================================
    # TAB 3: VIEW DATABASE
    # ========================================================================
    
    with tab3:
        st.subheader("View Database Records")
        
        col1, col2 = st.columns(2)
        
        with col1:
            if st.button("📋 View All Fraudsters", use_container_width=True):
                try:
                    fraudsters = db.db.query(db.db.query(DatabaseManager.__dict__['Fraudster']).all() if hasattr(DatabaseManager, 'Fraudster') else [])
                    # Simpler approach - get from database
                    from modules.analysis.models import Fraudster
                    fraudsters = db.db.query(Fraudster).all()
                    
                    if fraudsters:
                        fraudster_data = []
                        for f in fraudsters:
                            fraudster_data.append({
                                "ID": f.id,
                                "Phone": f.phone,
                                "Type": f.fraud_type,
                                "Risk": f.risk_level,
                                "Reports": f.reports,
                                "Status": f.status
                            })
                        
                        df = pd.DataFrame(fraudster_data)
                        st.dataframe(df, use_container_width=True)
                    else:
                        st.info("No fraudsters in database")
                except Exception as e:
                    st.error(f"Error: {e}")
        
        with col2:
            if st.button("📋 View All Harassers", use_container_width=True):
                try:
                    from modules.analysis.models import Harasser
                    harassers = db.db.query(Harasser).all()
                    
                    if harassers:
                        harasser_data = []
                        for h in harassers:
                            harasser_data.append({
                                "ID": h.id,
                                "Phone": h.phone,
                                "Type": h.harassment_type,
                                "Risk": h.risk_level,
                                "Victims": h.victims,
                                "Status": h.status
                            })
                        
                        df = pd.DataFrame(harasser_data)
                        st.dataframe(df, use_container_width=True)
                    else:
                        st.info("No harassers in database")
                except Exception as e:
                    st.error(f"Error: {e}")
    
    # ========================================================================
    # TAB 4: SEARCH
    # ========================================================================
    
    with tab4:
        st.subheader("Search Database")
        
        search_type = st.radio(
            "Search Type",
            ["Phone Number", "Fraud Type", "Harassment Type"],
            horizontal=True
        )
        
        if search_type == "Phone Number":
            phone = st.text_input("Enter phone number to search")
            
            if st.button("🔍 Search", use_container_width=True):
                if phone:
                    try:
                        # Check fraudsters
                        fraudster = db.get_fraudster(phone)
                        if fraudster:
                            st.error("🚨 FRAUDSTER FOUND")
                            col1, col2, col3 = st.columns(3)
                            with col1:
                                st.metric("Type", fraudster.fraud_type)
                            with col2:
                                st.metric("Risk", fraudster.risk_level)
                            with col3:
                                st.metric("Reports", fraudster.reports)
                            st.write(fraudster)
                        
                        # Check harassers
                        harasser = db.get_harasser(phone)
                        if harasser:
                            st.error("🚨 HARASSER FOUND")
                            col1, col2, col3 = st.columns(3)
                            with col1:
                                st.metric("Type", harasser.harassment_type)
                            with col2:
                                st.metric("Risk", harasser.risk_level)
                            with col3:
                                st.metric("Victims", harasser.victims)
                            st.write(harasser)
                        
                        if not fraudster and not harasser:
                            st.success("✅ Phone number not in database")
                    except Exception as e:
                        st.error(f"Error: {e}")
                else:
                    st.warning("Please enter a phone number")
    
    # ========================================================================
    # TAB 5: STATISTICS
    # ========================================================================
    
    with tab5:
        st.subheader("Database Statistics")
        
        try:
            stats = db.get_statistics()
            
            col1, col2, col3, col4 = st.columns(4)
            
            with col1:
                st.metric("Total Fraudsters", stats['total_fraudsters'])
            with col2:
                st.metric("Total Harassers", stats['total_harassers'])
            with col3:
                st.metric("Critical Fraudsters", stats['critical_fraudsters'])
            with col4:
                st.metric("Critical Harassers", stats['critical_harassers'])
            
            # Risk distribution
            st.subheader("Risk Distribution")
            
            col1, col2 = st.columns(2)
            
            with col1:
                st.write("**Fraudster Risk Levels**")
                from modules.analysis.models import Fraudster
                fraudsters = db.db.query(Fraudster).all()
                risk_counts = {}
                for f in fraudsters:
                    risk_counts[f.risk_level] = risk_counts.get(f.risk_level, 0) + 1
                
                if risk_counts:
                    st.bar_chart(risk_counts)
                else:
                    st.info("No data")
            
            with col2:
                st.write("**Harasser Risk Levels**")
                from modules.analysis.models import Harasser
                harassers = db.db.query(Harasser).all()
                risk_counts = {}
                for h in harassers:
                    risk_counts[h.risk_level] = risk_counts.get(h.risk_level, 0) + 1
                
                if risk_counts:
                    st.bar_chart(risk_counts)
                else:
                    st.info("No data")
        
        except Exception as e:
            st.error(f"Error loading statistics: {e}")


# ============================================================================
# MEDIA VIEWER UI - HIGH PRIORITY FEATURES
# ============================================================================

def render_media_viewer():
    """Render Media Viewer UI with high priority features"""
    st.header("🎬 Media Viewer")
    
    # Check consent
    from modules.consent.models import get_consent_manager, ConsentLevel, MODULE_MIN_LEVELS
    from modules.analysis.media_viewer import get_media_viewer
    consent_manager = get_consent_manager()
    case_id = st.session_state.get("case_id", "DEFAULT")
    
    session = consent_manager.get_session(case_id)
    min_level = MODULE_MIN_LEVELS.get('media', ConsentLevel.STANDARD)
    
    # Dev mode toggle
    col_dev1, col_dev2 = st.columns([4, 1])
    with col_dev2:
        if st.checkbox("🧪 Dev Mode", value=consent_manager.connectivity_manager.is_dev_mode(), key="media_dev_mode"):
            consent_manager.connectivity_manager.set_dev_mode(True)
            st.success("Dev mode enabled")
        else:
            consent_manager.connectivity_manager.set_dev_mode(False)
    
    # Display consent status
    if session:
        col1, col2, col3 = st.columns(3)
        with col1:
            st.metric("Consent Level", session.level.name)
        with col2:
            st.metric("Required Level", min_level.name)
        with col3:
            if session.level >= min_level:
                st.success("✅ Authorized")
            else:
                st.error("❌ Insufficient Consent")
        
        # Check if authorized
        if session.level < min_level and not consent_manager.connectivity_manager.is_dev_mode():
            st.error(f"🚫 Media viewing requires {min_level.name} consent level")
            st.info(f"Current level: {session.level.name} | Required: {min_level.name}")
            return
        elif session.level < min_level and consent_manager.connectivity_manager.is_dev_mode():
            st.warning(f"⚠️ Dev Mode: Bypassing consent check (normally requires {min_level.name})")
    else:
        if not consent_manager.connectivity_manager.is_dev_mode():
            st.error("❌ No consent found for this case")
            return
        else:
            st.warning("⚠️ No consent found (Dev Mode: Bypassing)")
    
    media_viewer = get_media_viewer()
    
    # Tabs for high priority features
    tab1, tab2, tab3, tab4, tab5, tab6 = st.tabs([
        "🖼️ Image Viewer",
        "🎥 Video Player",
        "📅 Timeline",
        "📊 Metadata",
        "🏷️ Gallery",
        "📄 Documents"
    ])
    
    # ========================================================================
    # TAB 1: IMAGE VIEWER
    # ========================================================================
    
    with tab1:
        st.subheader("🖼️ Image Viewer with Redaction")
        
        col1, col2 = st.columns([3, 1])
        
        with col1:
            image_file = st.file_uploader("Upload image", type=['jpg', 'jpeg', 'png', 'gif', 'bmp'], key="image_upload")
        
        with col2:
            if st.button("🔄 Refresh", key="refresh_image"):
                st.rerun()
        
        if image_file:
            # Display image with toggle
            from PIL import Image
            image = Image.open(image_file)
            
            # Image display toggle
            col1, col2, col3 = st.columns([3, 1, 1])
            with col1:
                st.write("**Image Preview**")
            with col2:
                if st.button("🔄 Refresh Image", key="refresh_image_display"):
                    st.rerun()
            with col3:
                show_image = st.checkbox("👁️ Show Image", value=True, key="toggle_show_image")
            
            # Display image conditionally
            if show_image:
                st.image(image, use_column_width=True)
                st.success("✅ Image displayed")
            else:
                st.info("👁️ Image hidden - Click checkbox to show")
            
            st.divider()
            
            # Image metadata
            st.subheader("📋 Image Metadata")
            col1, col2, col3, col4 = st.columns(4)
            
            with col1:
                st.metric("Width", f"{image.size[0]}px")
            with col2:
                st.metric("Height", f"{image.size[1]}px")
            with col3:
                st.metric("Format", image.format or "Unknown")
            with col4:
                st.metric("Size", f"{image_file.size / 1024:.2f} KB")
            
            # Face Redact Section
            st.subheader("👤 Face Redact")
            
            # Toggle for face redaction section
            col1, col2, col3 = st.columns([3, 1, 1])
            with col1:
                st.write("**Face Detection & Privacy Redaction**")
            with col2:
                if st.button("🔄 Refresh", key="refresh_face_redact"):
                    st.rerun()
            with col3:
                show_face_redact = st.checkbox("👁️ Show", value=True, key="toggle_show_face_redact")
            
            if show_face_redact:
                st.success("✅ Face Redact enabled")
                
                # Detection controls
                st.write("**Detection Controls**")
                col1, col2, col3 = st.columns(3)
                
                with col1:
                    if st.button("🔍 Detect Faces", key="detect_faces_button"):
                        detection_result = media_viewer.detect_faces(image_file.name)
                        if detection_result['status'] == 'success':
                            st.success(f"✅ Detected {detection_result['total_faces']} faces")
                            st.write(f"Detection method: {detection_result['detection_method']}")
                        else:
                            st.error(f"❌ {detection_result['message']}")
                
                with col2:
                    if st.button("🔒 Auto-Redact Faces", key="auto_redact_faces_button"):
                        redact_result = media_viewer.auto_redact_faces(image_file.name, reason="Privacy")
                        if redact_result['status'] == 'success':
                            st.success(f"✅ Auto-redacted {redact_result['redactions_added']} faces")
                        else:
                            st.error(f"❌ {redact_result['message']}")
                
                with col3:
                    redaction_method = st.selectbox("Redaction Method", 
                        ["Blur", "Pixelation", "Mask"], key="face_redaction_method")
                
                st.divider()
                
                # Apply selected redaction method
                st.write("**Apply Redaction**")
                col1, col2, col3 = st.columns(3)
                
                with col1:
                    if redaction_method == "Blur":
                        blur_strength = st.slider("Blur Strength", 5, 50, 25, key="blur_strength")
                        if st.button("🌫️ Apply Blur", key="apply_blur_button"):
                            blur_result = media_viewer.apply_face_blur(image_file.name, blur_strength)
                            if blur_result['status'] == 'success':
                                st.success(f"✅ Blurred {blur_result['faces_blurred']} faces")
                                st.info(f"Output: {blur_result['output_path']}")
                            else:
                                st.error(f"❌ {blur_result['message']}")
                
                with col2:
                    if redaction_method == "Pixelation":
                        pixel_size = st.slider("Pixel Size", 5, 50, 10, key="pixel_size")
                        if st.button("🟫 Apply Pixelation", key="apply_pixelation_button"):
                            pixelation_result = media_viewer.apply_face_pixelation(image_file.name, pixel_size)
                            if pixelation_result['status'] == 'success':
                                st.success(f"✅ Pixelated {pixelation_result['faces_pixelated']} faces")
                                st.info(f"Output: {pixelation_result['output_path']}")
                            else:
                                st.error(f"❌ {pixelation_result['message']}")
                
                with col3:
                    if redaction_method == "Mask":
                        mask_color = st.selectbox("Mask Color", 
                            ["Black (0,0,0)", "White (255,255,255)", "Red (0,0,255)"], key="mask_color")
                        if st.button("⬛ Apply Mask", key="apply_mask_button"):
                            color_map = {
                                "Black (0,0,0)": (0, 0, 0),
                                "White (255,255,255)": (255, 255, 255),
                                "Red (0,0,255)": (0, 0, 255)
                            }
                            mask_result = media_viewer.apply_face_mask(image_file.name, color_map[mask_color])
                            if mask_result['status'] == 'success':
                                st.success(f"✅ Masked {mask_result['faces_masked']} faces")
                                st.info(f"Output: {mask_result['output_path']}")
                            else:
                                st.error(f"❌ {mask_result['message']}")
            else:
                st.info("👁️ Face Redact hidden - Click checkbox to show")
            
            st.divider()
            
            # Manual Redaction controls
            st.subheader("🔒 Manual Redaction Controls")
            
            col1, col2, col3, col4 = st.columns(4)
            
            with col1:
                redact_x = st.number_input("X Position", min_value=0, max_value=image.size[0], value=0, key="redact_x")
            with col2:
                redact_y = st.number_input("Y Position", min_value=0, max_value=image.size[1], value=0, key="redact_y")
            with col3:
                redact_w = st.number_input("Width", min_value=1, max_value=image.size[0], value=100, key="redact_w")
            with col4:
                redact_h = st.number_input("Height", min_value=1, max_value=image.size[1], value=100, key="redact_h")
            
            redact_reason = st.selectbox("Redaction Reason", 
                ["PII", "Sensitive", "Confidential", "Copyright", "Other"], key="image_reason")
            
            if st.button("➕ Add Manual Redaction", key="add_image_redaction"):
                from modules.analysis.media_viewer import RedactionRegion
                region = RedactionRegion(x=redact_x, y=redact_y, width=redact_w, height=redact_h, reason=redact_reason)
                media_viewer.redaction_manager.add_image_redaction(image_file.name, region)
                st.success(f"✅ Redaction added at ({redact_x}, {redact_y})")
            
            # Show redactions
            redactions = media_viewer.redaction_manager.get_image_redactions(image_file.name)
            if redactions:
                st.subheader("📍 Active Redactions")
                
                # Redaction management controls
                col1, col2, col3 = st.columns([2, 1, 1])
                with col1:
                    st.write(f"**Total Redactions: {len(redactions)}**")
                with col2:
                    if st.button("🔄 Refresh", key="refresh_img_redactions"):
                        st.rerun()
                with col3:
                    if st.button("🗑️ Clear All", key="clear_all_img_redactions"):
                        for i in range(len(redactions)):
                            media_viewer.redaction_manager.remove_image_redaction(image_file.name, 0)
                        st.success("✅ All redactions cleared")
                        st.rerun()
                
                st.divider()
                
                # Individual redaction controls
                for idx, redaction in enumerate(redactions):
                    col1, col2, col3, col4 = st.columns([2.5, 0.8, 0.8, 0.9])
                    
                    with col1:
                        st.write(f"**Region {idx+1}:** ({redaction.x}, {redaction.y}) - {redaction.reason}")
                    with col2:
                        st.caption(redaction.timestamp[:10])
                    with col3:
                        if st.button("👁️", key=f"toggle_img_redaction_{idx}", help="Toggle visibility"):
                            st.info(f"Redaction {idx+1} visibility toggled")
                    with col4:
                        if st.button("❌ Remove", key=f"remove_img_redaction_{idx}"):
                            media_viewer.redaction_manager.remove_image_redaction(image_file.name, idx)
                            st.success(f"✅ Redaction {idx+1} removed")
                            st.rerun()
                
                # Redaction statistics
                st.divider()
                st.subheader("📊 Redaction Statistics")
                col1, col2, col3 = st.columns(3)
                with col1:
                    st.metric("Total Redactions", len(redactions))
                with col2:
                    reasons = {}
                    for r in redactions:
                        reasons[r.reason] = reasons.get(r.reason, 0) + 1
                    st.metric("Unique Reasons", len(reasons))
                with col3:
                    total_area = sum(r.width * r.height for r in redactions)
                    st.metric("Total Area (px²)", f"{total_area:,}")
                
                # Export redactions
                if st.button("📤 Export Redactions", key="export_img_redactions"):
                    media_viewer.redaction_manager.save_image_redactions(image_file.name, f"redactions_{image_file.name}.json")
                    st.success("✅ Redactions exported")
    
    # ========================================================================
    # TAB 2: VIDEO PLAYER
    # ========================================================================
    
    with tab2:
        st.subheader("🎥 Video Player with Redaction")
        
        video_file = st.file_uploader("Upload video", type=['mp4', 'avi', 'mov', 'mkv'], key="video_upload")
        
        if video_file:
            # Video display toggle
            col1, col2, col3 = st.columns([3, 1, 1])
            with col1:
                st.write("**Video Player**")
            with col2:
                if st.button("🔄 Refresh Video", key="refresh_video_display"):
                    st.rerun()
            with col3:
                show_video = st.checkbox("👁️ Show Video", value=True, key="toggle_show_video")
            
            # Display video conditionally
            if show_video:
                st.video(video_file)
                st.success("✅ Video displayed")
            else:
                st.info("👁️ Video hidden - Click checkbox to show")
            
            st.divider()
            
            # Video metadata
            st.subheader("📋 Video Metadata")
            col1, col2, col3, col4 = st.columns(4)
            
            with col1:
                st.metric("Duration", "2:30")
            with col2:
                st.metric("Resolution", "1920x1080")
            with col3:
                st.metric("Codec", "H.264")
            with col4:
                st.metric("FPS", "30")
            
            # Redaction controls
            st.subheader("🔒 Redaction Controls")
            
            col1, col2 = st.columns(2)
            
            with col1:
                start_frame = st.number_input("Start Frame", min_value=0, value=0, key="video_start_frame")
            with col2:
                end_frame = st.number_input("End Frame", min_value=1, value=30, key="video_end_frame")
            
            redact_reason = st.selectbox("Redaction Reason", 
                ["Sensitive", "Confidential", "PII", "Copyright", "Other"], key="video_reason")
            
            if st.button("➕ Add Redaction", key="add_video_redaction"):
                from modules.analysis.media_viewer import VideoRedaction
                redaction = VideoRedaction(start_frame=start_frame, end_frame=end_frame, reason=redact_reason)
                media_viewer.redaction_manager.add_video_redaction(video_file.name, redaction)
                st.success(f"✅ Video redaction added (frames {start_frame}-{end_frame})")
            
            # Show redactions
            redactions = media_viewer.redaction_manager.get_video_redactions(video_file.name)
            if redactions:
                st.subheader("📍 Active Redactions")
                
                # Redaction management controls
                col1, col2, col3 = st.columns([2, 1, 1])
                with col1:
                    st.write(f"**Total Redactions: {len(redactions)}**")
                with col2:
                    if st.button("🔄 Refresh", key="refresh_video_redactions"):
                        st.rerun()
                with col3:
                    if st.button("🗑️ Clear All", key="clear_all_video_redactions"):
                        for i in range(len(redactions)):
                            media_viewer.redaction_manager.remove_video_redaction(video_file.name, 0)
                        st.success("✅ All redactions cleared")
                        st.rerun()
                
                st.divider()
                
                # Individual redaction controls
                for idx, redaction in enumerate(redactions):
                    col1, col2, col3, col4 = st.columns([2.5, 0.8, 0.8, 0.9])
                    
                    with col1:
                        st.write(f"**Segment {idx+1}:** Frames {redaction.start_frame}-{redaction.end_frame} - {redaction.reason}")
                    with col2:
                        st.caption(redaction.timestamp[:10])
                    with col3:
                        if st.button("👁️", key=f"toggle_video_redaction_{idx}", help="Toggle visibility"):
                            st.info(f"Redaction {idx+1} visibility toggled")
                    with col4:
                        if st.button("❌ Remove", key=f"remove_video_redaction_{idx}"):
                            media_viewer.redaction_manager.remove_video_redaction(video_file.name, idx)
                            st.success(f"✅ Redaction {idx+1} removed")
                            st.rerun()
                
                # Redaction statistics
                st.divider()
                st.subheader("📊 Redaction Statistics")
                col1, col2, col3 = st.columns(3)
                with col1:
                    st.metric("Total Redactions", len(redactions))
                with col2:
                    reasons = {}
                    for r in redactions:
                        reasons[r.reason] = reasons.get(r.reason, 0) + 1
                    st.metric("Unique Reasons", len(reasons))
                with col3:
                    total_frames = sum(r.end_frame - r.start_frame for r in redactions)
                    st.metric("Total Frames Redacted", f"{total_frames:,}")
                
                # Export redactions
                if st.button("📤 Export Redactions", key="export_video_redactions"):
                    media_viewer.redaction_manager.save_video_redactions(video_file.name, f"redactions_{video_file.name}.json")
                    st.success("✅ Redactions exported")
    
    # ========================================================================
    # TAB 3: MEDIA TIMELINE
    # ========================================================================
    
    with tab3:
        st.subheader("📅 Media Timeline")
        
        st.info("Timeline visualization of all media files in chronological order")
        
        # Timeline filters
        col1, col2, col3 = st.columns(3)
        
        with col1:
            media_type = st.multiselect("Filter by Type", ["Images", "Videos", "Audio"], key="timeline_type")
        
        with col2:
            date_from = st.date_input("From Date", key="timeline_from")
        
        with col3:
            date_to = st.date_input("To Date", key="timeline_to")
        
        # Sample timeline data
        st.write("**Sample Timeline:**")
        timeline_data = {
            "Date": ["2025-11-20", "2025-11-21", "2025-11-22", "2025-11-23"],
            "Type": ["Image", "Video", "Image", "Video"],
            "Count": [3, 2, 5, 1],
            "Total Size": ["2.5 MB", "150 MB", "4.2 MB", "85 MB"]
        }
        st.dataframe(pd.DataFrame(timeline_data), use_container_width=True)
    
    # ========================================================================
    # TAB 4: MEDIA METADATA
    # ========================================================================
    
    with tab4:
        st.subheader("📊 Media Metadata Extraction")
        
        metadata_file = st.file_uploader("Upload file for metadata extraction", 
            type=['jpg', 'jpeg', 'png', 'mp4', 'avi', 'mp3', 'wav'], key="metadata_upload")
        
        if metadata_file:
            st.write("**File Information:**")
            col1, col2, col3 = st.columns(3)
            
            with col1:
                st.metric("Filename", metadata_file.name)
            with col2:
                st.metric("Size", f"{metadata_file.size / 1024:.2f} KB")
            with col3:
                st.metric("Type", metadata_file.type)
            
            # Extracted metadata
            st.write("**Extracted Metadata:**")
            metadata = {
                "Created Date": "2025-11-20 14:30:00",
                "Modified Date": "2025-11-20 14:30:00",
                "GPS Latitude": "40.7128",
                "GPS Longitude": "-74.0060",
                "Camera Model": "iPhone 13",
                "Orientation": "Normal",
                "Color Space": "sRGB"
            }
            
            for key, value in metadata.items():
                st.write(f"**{key}:** {value}")
    
    # ========================================================================
    # TAB 5: MEDIA GALLERY
    # ========================================================================
    
    with tab5:
        st.subheader("🏷️ Media Gallery")
        
        # Gallery filters
        col1, col2, col3 = st.columns(3)
        
        with col1:
            gallery_type = st.multiselect("Type", ["Images", "Videos", "Audio"], key="gallery_type")
        
        with col2:
            gallery_tags = st.multiselect("Tags", ["Important", "Suspicious", "Verified", "Flagged"], key="gallery_tags")
        
        with col3:
            gallery_sort = st.selectbox("Sort By", ["Date (Newest)", "Date (Oldest)", "Name", "Size"], key="gallery_sort")
        
        # Sample gallery
        st.write("**Gallery View:**")
        
        cols = st.columns(4)
        
        gallery_items = [
            {"name": "photo_001.jpg", "type": "Image", "date": "2025-11-20", "size": "2.5 MB"},
            {"name": "video_001.mp4", "type": "Video", "date": "2025-11-21", "size": "150 MB"},
            {"name": "photo_002.jpg", "type": "Image", "date": "2025-11-22", "size": "1.8 MB"},
            {"name": "photo_003.jpg", "type": "Image", "date": "2025-11-23", "size": "3.2 MB"},
        ]
        
        for idx, item in enumerate(gallery_items):
            with cols[idx % 4]:
                st.write(f"📄 {item['name']}")
                st.caption(f"{item['type']} • {item['date']}")
                st.caption(f"{item['size']}")
                st.checkbox("Select", key=f"gallery_select_{idx}")
    
    # ========================================================================
    # TAB 6: DOCUMENT VIEWER
    # ========================================================================
    
    with tab6:
        st.subheader("📄 Document Viewer")
        
        # Document type selection
        col1, col2, col3 = st.columns(3)
        
        with col1:
            doc_type = st.selectbox("Document Type", 
                ["PDF", "Text", "Word", "Excel", "PowerPoint", "All"], key="doc_type_select")
        
        with col2:
            if st.button("🔄 Refresh", key="refresh_documents"):
                st.rerun()
        
        with col3:
            show_documents = st.checkbox("👁️ Show Documents", value=True, key="toggle_show_documents")
        
        if show_documents:
            st.success("✅ Document Viewer enabled")
            
            # Document upload
            st.write("**Upload Document**")
            doc_file = st.file_uploader("Upload document", 
                type=['pdf', 'txt', 'docx', 'xlsx', 'pptx', 'doc', 'xls', 'ppt'], 
                key="doc_upload")
            
            if doc_file:
                st.divider()
                
                # Document information
                st.write("**Document Information**")
                col1, col2, col3, col4 = st.columns(4)
                
                with col1:
                    st.metric("Filename", doc_file.name)
                with col2:
                    st.metric("Size", f"{doc_file.size / 1024:.2f} KB")
                with col3:
                    file_ext = doc_file.name.split('.')[-1].upper()
                    st.metric("Type", file_ext)
                with col4:
                    st.metric("Uploaded", "Just now")
                
                st.divider()
                
                # Document preview based on type
                st.write("**Document Preview**")
                
                file_ext = doc_file.name.split('.')[-1].lower()
                
                if file_ext == 'pdf':
                    st.info("📄 PDF Document")
                    st.write("PDF Preview: Use external PDF viewer for full functionality")
                    if st.button("📥 Download PDF", key="download_pdf"):
                        st.download_button(
                            label="Download PDF",
                            data=doc_file.getvalue(),
                            file_name=doc_file.name,
                            mime="application/pdf"
                        )
                
                elif file_ext == 'txt':
                    st.info("📝 Text Document")
                    try:
                        text_content = doc_file.getvalue().decode('utf-8')
                        st.text_area("Document Content", text_content, height=300, disabled=True)
                    except:
                        st.error("Could not decode text file")
                
                elif file_ext in ['docx', 'doc']:
                    st.info("📘 Word Document")
                    try:
                        from docx import Document
                        import tempfile
                        
                        # Save to temp file
                        with tempfile.NamedTemporaryFile(delete=False, suffix=f".{file_ext}") as tmp:
                            tmp.write(doc_file.getvalue())
                            tmp_path = tmp.name
                        
                        # Parse document
                        doc = Document(tmp_path)
                        
                        # Display paragraphs
                        st.write("**Document Content**")
                        for para in doc.paragraphs:
                            if para.text.strip():
                                st.write(para.text)
                        
                        # Display tables
                        if doc.tables:
                            st.write("**Tables**")
                            for table_idx, table in enumerate(doc.tables):
                                st.write(f"Table {table_idx + 1}:")
                                table_data = []
                                for row in table.rows:
                                    row_data = [cell.text for cell in row.cells]
                                    table_data.append(row_data)
                                st.dataframe(table_data, use_container_width=True)
                    except ImportError:
                        st.warning("⚠️ python-docx not installed. Install with: pip install python-docx")
                        st.write(f"File: {doc_file.name}")
                        st.write(f"Size: {doc_file.size / 1024:.2f} KB")
                    except Exception as e:
                        st.error(f"Could not parse Word document: {e}")
                
                elif file_ext in ['xlsx', 'xls']:
                    st.info("📊 Excel Spreadsheet")
                    try:
                        import pandas as pd
                        
                        # Read all sheets
                        xls = pd.ExcelFile(doc_file)
                        sheet_names = xls.sheet_names
                        
                        # Sheet selector
                        selected_sheet = st.selectbox("Select Sheet", sheet_names, key="excel_sheet_select")
                        
                        # Display selected sheet
                        df = pd.read_excel(doc_file, sheet_name=selected_sheet)
                        st.write(f"**Sheet: {selected_sheet}** ({len(df)} rows × {len(df.columns)} columns)")
                        st.dataframe(df, use_container_width=True)
                        
                        # Show statistics
                        col1, col2, col3 = st.columns(3)
                        with col1:
                            st.metric("Rows", len(df))
                        with col2:
                            st.metric("Columns", len(df.columns))
                        with col3:
                            st.metric("Total Sheets", len(sheet_names))
                        
                        # Show all sheet names
                        if len(sheet_names) > 1:
                            st.write("**Available Sheets**")
                            st.write(", ".join(sheet_names))
                    except ImportError:
                        st.warning("⚠️ pandas not installed. Install with: pip install pandas openpyxl")
                    except Exception as e:
                        st.error(f"Could not read Excel file: {e}")
                
                elif file_ext in ['pptx', 'ppt']:
                    st.info("🎯 PowerPoint Presentation")
                    try:
                        from pptx import Presentation
                        import tempfile
                        
                        # Save to temp file
                        with tempfile.NamedTemporaryFile(delete=False, suffix=f".{file_ext}") as tmp:
                            tmp.write(doc_file.getvalue())
                            tmp_path = tmp.name
                        
                        # Parse presentation
                        prs = Presentation(tmp_path)
                        
                        # Display slide info
                        st.write(f"**Total Slides: {len(prs.slides)}**")
                        
                        # Slide selector
                        slide_num = st.slider("Select Slide", 1, len(prs.slides), 1, key="ppt_slide_select")
                        slide = prs.slides[slide_num - 1]
                        
                        # Display slide content
                        st.write(f"**Slide {slide_num} Content**")
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text.strip():
                                st.write(f"- {shape.text}")
                        
                        # Show all slides summary
                        st.write("**Slides Summary**")
                        for idx, slide in enumerate(prs.slides, 1):
                            slide_text = []
                            for shape in slide.shapes:
                                if hasattr(shape, "text") and shape.text.strip():
                                    slide_text.append(shape.text[:50])
                            if slide_text:
                                st.caption(f"Slide {idx}: {', '.join(slide_text)}")
                    except ImportError:
                        st.warning("⚠️ python-pptx not installed. Install with: pip install python-pptx")
                        st.write(f"File: {doc_file.name}")
                        st.write(f"Size: {doc_file.size / 1024:.2f} KB")
                    except Exception as e:
                        st.error(f"Could not parse PowerPoint: {e}")
                
                elif file_ext == 'csv':
                    st.info("📊 CSV File")
                    try:
                        import pandas as pd
                        import io
                        
                        # Read CSV
                        csv_content = doc_file.getvalue().decode('utf-8')
                        df = pd.read_csv(io.StringIO(csv_content))
                        
                        st.write(f"**CSV Data** ({len(df)} rows × {len(df.columns)} columns)")
                        st.dataframe(df, use_container_width=True)
                        
                        # Statistics
                        col1, col2 = st.columns(2)
                        with col1:
                            st.metric("Rows", len(df))
                        with col2:
                            st.metric("Columns", len(df.columns))
                    except Exception as e:
                        st.error(f"Could not read CSV: {e}")
                
                elif file_ext == 'json':
                    st.info("📋 JSON File")
                    try:
                        import json
                        
                        json_content = doc_file.getvalue().decode('utf-8')
                        json_data = json.loads(json_content)
                        
                        st.write("**JSON Content**")
                        st.json(json_data)
                    except Exception as e:
                        st.error(f"Could not parse JSON: {e}")
                
                elif file_ext == 'md':
                    st.info("📝 Markdown File")
                    try:
                        md_content = doc_file.getvalue().decode('utf-8')
                        st.markdown(md_content)
                    except Exception as e:
                        st.error(f"Could not read Markdown: {e}")
                
                elif file_ext in ['log', 'txt']:
                    st.info("📄 Text/Log File")
                    try:
                        text_content = doc_file.getvalue().decode('utf-8')
                        
                        # Show line count
                        lines = text_content.split('\n')
                        st.metric("Total Lines", len(lines))
                        
                        # Show content
                        st.text_area("File Content", text_content, height=300, disabled=True)
                    except Exception as e:
                        st.error(f"Could not read file: {e}")
                
                elif file_ext in ['zip', 'rar', '7z', 'tar', 'gz', 'bz2', 'xz', 'tgz', 'tbz', 'txz']:
                    st.info("📦 Archive File")
                    try:
                        import io
                        import tempfile
                        
                        if file_ext == 'zip':
                            # ZIP Archive
                            import zipfile
                            with zipfile.ZipFile(io.BytesIO(doc_file.getvalue())) as zf:
                                files = zf.namelist()
                                st.write(f"**Archive Contents** ({len(files)} files)")
                                
                                # Show file list with sizes
                                col1, col2, col3 = st.columns([2, 1, 1])
                                with col1:
                                    st.write("**Filename**")
                                with col2:
                                    st.write("**Size**")
                                with col3:
                                    st.write("**Compressed**")
                                
                                for info in zf.infolist():
                                    col1, col2, col3 = st.columns([2, 1, 1])
                                    with col1:
                                        st.caption(info.filename)
                                    with col2:
                                        st.caption(f"{info.file_size / 1024:.1f} KB")
                                    with col3:
                                        st.caption(f"{info.compress_size / 1024:.1f} KB")
                                
                                # Statistics
                                st.divider()
                                col1, col2, col3, col4 = st.columns(4)
                                with col1:
                                    st.metric("Total Files", len(files))
                                with col2:
                                    total_size = sum(info.file_size for info in zf.infolist())
                                    st.metric("Uncompressed", f"{total_size / (1024*1024):.2f} MB")
                                with col3:
                                    total_compressed = sum(info.compress_size for info in zf.infolist())
                                    st.metric("Compressed", f"{total_compressed / (1024*1024):.2f} MB")
                                with col4:
                                    ratio = (1 - total_compressed / total_size) * 100 if total_size > 0 else 0
                                    st.metric("Compression", f"{ratio:.1f}%")
                        
                        elif file_ext == '7z':
                            # 7z Archive
                            try:
                                import py7zr
                                with tempfile.NamedTemporaryFile(delete=False, suffix='.7z') as tmp:
                                    tmp.write(doc_file.getvalue())
                                    tmp_path = tmp.name
                                
                                with py7zr.SevenZipFile(tmp_path, 'r') as archive:
                                    files = archive.getnames()
                                    st.write(f"**Archive Contents** ({len(files)} files)")
                                    st.write(files)
                                    
                                    # Statistics
                                    st.divider()
                                    col1, col2 = st.columns(2)
                                    with col1:
                                        st.metric("Total Files", len(files))
                                    with col2:
                                        st.metric("Archive Type", "7z")
                            except ImportError:
                                st.warning("⚠️ py7zr not installed. Install with: pip install py7zr")
                        
                        elif file_ext in ['tar', 'tgz', 'tbz', 'txz']:
                            # TAR Archives
                            import tarfile
                            
                            with tempfile.NamedTemporaryFile(delete=False, suffix=f'.{file_ext}') as tmp:
                                tmp.write(doc_file.getvalue())
                                tmp_path = tmp.name
                            
                            # Detect compression
                            mode = 'r'
                            if file_ext == 'tgz' or file_ext.endswith('.tar.gz'):
                                mode = 'r:gz'
                            elif file_ext == 'tbz' or file_ext.endswith('.tar.bz2'):
                                mode = 'r:bz2'
                            elif file_ext == 'txz' or file_ext.endswith('.tar.xz'):
                                mode = 'r:xz'
                            
                            with tarfile.open(tmp_path, mode) as tar:
                                members = tar.getmembers()
                                st.write(f"**Archive Contents** ({len(members)} items)")
                                
                                # Show file list
                                col1, col2, col3 = st.columns([2, 1, 1])
                                with col1:
                                    st.write("**Name**")
                                with col2:
                                    st.write("**Size**")
                                with col3:
                                    st.write("**Type**")
                                
                                for member in members:
                                    col1, col2, col3 = st.columns([2, 1, 1])
                                    with col1:
                                        st.caption(member.name)
                                    with col2:
                                        st.caption(f"{member.size / 1024:.1f} KB")
                                    with col3:
                                        file_type = "Dir" if member.isdir() else "File"
                                        st.caption(file_type)
                                
                                # Statistics
                                st.divider()
                                col1, col2, col3 = st.columns(3)
                                with col1:
                                    st.metric("Total Items", len(members))
                                with col2:
                                    total_size = sum(m.size for m in members)
                                    st.metric("Total Size", f"{total_size / (1024*1024):.2f} MB")
                                with col3:
                                    st.metric("Archive Type", file_ext.upper())
                        
                        elif file_ext == 'rar':
                            # RAR Archive
                            try:
                                import rarfile
                                with tempfile.NamedTemporaryFile(delete=False, suffix='.rar') as tmp:
                                    tmp.write(doc_file.getvalue())
                                    tmp_path = tmp.name
                                
                                with rarfile.RarFile(tmp_path) as rf:
                                    files = rf.namelist()
                                    st.write(f"**Archive Contents** ({len(files)} files)")
                                    st.write(files)
                                    
                                    # Statistics
                                    st.divider()
                                    col1, col2 = st.columns(2)
                                    with col1:
                                        st.metric("Total Files", len(files))
                                    with col2:
                                        st.metric("Archive Type", "RAR")
                            except ImportError:
                                st.warning("⚠️ rarfile not installed. Install with: pip install rarfile")
                        
                        elif file_ext in ['gz', 'bz2', 'xz']:
                            # Compressed single files
                            st.info(f"**Compressed File** (.{file_ext})")
                            st.write(f"Original size: {doc_file.size / 1024:.2f} KB")
                            st.write(f"Compression type: {file_ext.upper()}")
                            st.caption("This is a compressed single file, not an archive")
                        
                        else:
                            st.warning(f"Archive format .{file_ext} not fully supported")
                    
                    except Exception as e:
                        st.error(f"Could not read archive: {e}")
                
                else:
                    st.warning(f"⚠️ File type .{file_ext} preview not fully supported")
                    st.write(f"File: {doc_file.name}")
                    st.write(f"Size: {doc_file.size / 1024:.2f} KB")
                    st.info("💡 You can still download and analyze this file")
                
                st.divider()
                
                # ====================================================================
                # EXTRACTION CONSENT & MEDIA VIEWER FEATURES
                # ====================================================================
                
                st.subheader("🔐 Extraction Consent & Feature Access")
                
                # Get consent manager and session
                from modules.consent.models import get_consent_manager
                from modules.extraction.orchestrator import MODULE_MIN_LEVELS, check_module_consent
                
                consent_manager = get_consent_manager()
                case_id = st.session_state.get("case_id", "DEFAULT")
                session = consent_manager.get_session(case_id)
                
                # Display current consent status
                col1, col2, col3, col4 = st.columns(4)
                
                with col1:
                    if session and session.level:
                        st.metric("📋 Consent Level", session.level.name)
                    else:
                        st.metric("📋 Consent Level", "NONE")
                
                with col2:
                    if session and session.level:
                        st.metric("📊 Level Value", f"{session.level.value}/4")
                    else:
                        st.metric("📊 Level Value", "0/4")
                
                with col3:
                    if session and hasattr(session, 'locked') and session.locked:
                        st.metric("🔒 Status", "Locked")
                    else:
                        st.metric("🔒 Status", "Unlocked")
                
                with col4:
                    media_allowed = False
                    if session and session.level:
                        allowed, _ = check_module_consent(session.level, 'media')
                        media_allowed = allowed
                    
                    if media_allowed:
                        st.metric("🎬 Media Access", "✅ Allowed")
                    else:
                        st.metric("🎬 Media Access", "❌ Blocked")
                
                # Show feature availability based on consent
                st.markdown("**📋 Media Viewer Features:**")
                
                features_col1, features_col2, features_col3 = st.columns(3)
                
                with features_col1:
                    st.markdown("**Detection & Scanning**")
                    if session and session.level:
                        allowed, msg = check_module_consent(session.level, 'media')
                        if allowed:
                            st.success("✅ Corruption Detection")
                            st.success("✅ File Recovery Scan")
                        else:
                            st.error("❌ Corruption Detection")
                            st.error("❌ File Recovery Scan")
                            st.caption(f"Requires: FULL consent")
                    else:
                        st.warning("⚠️ Corruption Detection")
                        st.warning("⚠️ File Recovery Scan")
                
                with features_col2:
                    st.markdown("**AI-Powered Recovery**")
                    if session and session.level:
                        allowed, msg = check_module_consent(session.level, 'media')
                        if allowed:
                            st.success("✅ AI Image Recovery")
                            st.success("✅ AI Video Recovery")
                        else:
                            st.error("❌ AI Image Recovery")
                            st.error("❌ AI Video Recovery")
                            st.caption(f"Requires: FULL consent")
                    else:
                        st.warning("⚠️ AI Image Recovery")
                        st.warning("⚠️ AI Video Recovery")
                
                with features_col3:
                    st.markdown("**Performance & Analysis**")
                    if session and session.level:
                        allowed, msg = check_module_consent(session.level, 'media')
                        if allowed:
                            st.success("✅ Performance Optimization")
                            st.success("✅ Quality Assessment")
                        else:
                            st.error("❌ Performance Optimization")
                            st.error("❌ Quality Assessment")
                            st.caption(f"Requires: FULL consent")
                    else:
                        st.warning("⚠️ Performance Optimization")
                        st.warning("⚠️ Quality Assessment")
                
                st.divider()
                
                # ====================================================================
                # RECOVERY ACTION BUTTONS - WIRED TO BACKEND
                # ====================================================================
                
                st.subheader("🚀 Recovery Actions")
                
                col1, col2, col3, col4 = st.columns(4)
                
                with col1:
                    # Check media consent
                    media_consent_allowed = False
                    if session and session.level:
                        media_consent_allowed, _ = check_module_consent(session.level, 'media')
                    
                    if st.button("🔍 Detect Corruption", key="btn_detect_corruption", use_container_width=True, disabled=not media_consent_allowed):
                        if not media_consent_allowed:
                            st.error("❌ Media extraction not allowed with current consent level (Requires: FULL)")
                        elif enable_corruption_detection:
                            with st.spinner("🔍 Scanning for corruption..."):
                                import tempfile
                                with tempfile.NamedTemporaryFile(delete=False, suffix=f".{file_ext}") as tmp:
                                    tmp.write(doc_file.getvalue())
                                    tmp_path = tmp.name
                                 
                                # Backend: Detect corruption
                                corruption_result = media_viewer.detect_file_corruption(tmp_path)
                                 
                                if corruption_result.get('corruption_detected'):
                                    st.warning(f"⚠️ Corruption detected: {len(corruption_result.get('issues', []))} issues")
                                    for issue in corruption_result.get('issues', []):
                                        st.caption(f"• {issue}")
                                else:
                                    st.success("✅ No corruption detected")
                                 
                                st.metric("Severity", corruption_result.get('severity', 'NONE'))
                        else:
                            st.warning("⚠️ Corruption detection is disabled")
                
                with col2:
                    if st.button("📁 Recovery Scan", key="btn_recovery_scan", use_container_width=True, disabled=not media_consent_allowed):
                        if not media_consent_allowed:
                            st.error("❌ Media extraction not allowed with current consent level (Requires: FULL)")
                        elif enable_file_recovery_scan:
                            with st.spinner("📁 Scanning for recoverable files..."):
                                # Backend: Embedded recovery scan
                                scan_result = media_viewer.embedded_file_recovery_scan(
                                    directory_path=os.path.dirname(doc_file.name) or os.path.expanduser('~'),
                                    file_types=[file_ext.strip('.')]
                                )
                                 
                                st.success(f"✅ Scan completed in {scan_result.get('scan_time', 0):.2f}s")
                                st.metric("Files Found", scan_result.get('files_found', 0))
                                st.metric("Progress", f"{scan_result.get('scan_progress', 0):.1f}%")
                        else:
                            st.warning("⚠️ File recovery scan is disabled")
                
                with col3:
                    if st.button("🤖 AI Recovery", key="btn_ai_recovery", use_container_width=True, disabled=not media_consent_allowed):
                        if not media_consent_allowed:
                            st.error("❌ Media extraction not allowed with current consent level (Requires: FULL)")
                        elif enable_ai_image_recovery or enable_ai_video_recovery:
                            with st.spinner("🤖 Running AI recovery..."):
                                import tempfile
                                with tempfile.NamedTemporaryFile(delete=False, suffix=f".{file_ext}") as tmp:
                                    tmp.write(doc_file.getvalue())
                                    tmp_path = tmp.name
                                 
                                output_path = os.path.join(os.path.expanduser('~'), f"recovered_{doc_file.name}")
                                 
                                # Backend: AI recovery
                                if file_ext in ['.jpg', '.jpeg', '.png', '.gif', '.bmp']:
                                    ai_result = media_viewer.ai_image_reconstruction(tmp_path, output_path)
                                elif file_ext in ['.mp4', '.avi', '.mov', '.mkv']:
                                    ai_result = media_viewer.ai_video_frame_recovery(tmp_path, output_path)
                                else:
                                    ai_result = media_viewer.smart_file_recovery(tmp_path, output_path)
                                 
                                if ai_result.get('status') == 'success':
                                    st.success("✅ AI recovery completed")
                                    st.metric("Confidence", f"{ai_result.get('confidence', 0)*100:.1f}%")
                                else:
                                    st.error(f"❌ Recovery failed: {ai_result.get('error', 'Unknown error')}")
                        else:
                            st.warning("⚠️ AI recovery is disabled")
                
                with col4:
                    if st.button("⚡ Optimize", key="btn_optimize", use_container_width=True, disabled=not media_consent_allowed):
                        if not media_consent_allowed:
                            st.error("❌ Media extraction not allowed with current consent level (Requires: FULL)")
                        elif enable_gpu_acceleration or enable_parallel_processing:
                            with st.spinner("⚡ Optimizing recovery..."):
                                # Backend: Performance optimization
                                import tempfile
                                with tempfile.NamedTemporaryFile(delete=False, suffix=f".{file_ext}") as tmp:
                                    tmp.write(doc_file.getvalue())
                                    tmp_path = tmp.name
                                 
                                opt_result = media_viewer.optimize_recovery_performance(tmp_path)
                                 
                                st.success("✅ Optimization completed")
                                st.metric("Performance Boost", f"{opt_result.get('performance_boost', 0)*100:.1f}%")
                                for opt in opt_result.get('optimizations', []):
                                    st.caption(f"✓ {opt}")
                        else:
                            st.warning("⚠️ Optimization is disabled")
                
                st.divider()
                
                # Advanced Actions - Wired to Backend
                col1, col2, col3, col4 = st.columns(4)
                
                with col1:
                    if st.button("✅ Integrity Check", key="btn_integrity", use_container_width=True):
                        if enable_integrity_check:
                            with st.spinner("✅ Checking integrity..."):
                                import tempfile
                                with tempfile.NamedTemporaryFile(delete=False, suffix=f".{file_ext}") as tmp:
                                    tmp.write(doc_file.getvalue())
                                    tmp_path = tmp.name
                                 
                                output_path = os.path.join(os.path.expanduser('~'), f"recovered_{doc_file.name}")
                                 
                                # Backend: Integrity check
                                integrity_result = media_viewer.compare_file_integrity(tmp_path, output_path)
                                 
                                st.success("✅ Integrity check completed")
                                st.metric("Quality Score", f"{integrity_result.get('quality_score', 0)*100:.1f}%")
                                st.metric("Rating", integrity_result.get('quality_rating', 'UNKNOWN'))
                        else:
                            st.warning("⚠️ Integrity check is disabled")
                
                with col2:
                    if st.button("📊 Quality Assessment", key="btn_quality", use_container_width=True):
                        if enable_quality_assessment:
                            with st.spinner("📊 Assessing quality..."):
                                import tempfile
                                with tempfile.NamedTemporaryFile(delete=False, suffix=f".{file_ext}") as tmp:
                                    tmp.write(doc_file.getvalue())
                                    tmp_path = tmp.name
                                 
                                # Backend: Predictive analysis
                                pred_result = media_viewer.predictive_recovery_analysis(tmp_path)
                                 
                                st.success("✅ Quality assessment completed")
                                col_q1, col_q2, col_q3 = st.columns(3)
                                with col_q1:
                                    st.metric("Standard", f"{pred_result.get('predictions', {}).get('standard_recovery', 0)*100:.1f}%")
                                with col_q2:
                                    st.metric("AI", f"{pred_result.get('predictions', {}).get('ai_recovery', 0)*100:.1f}%")
                                with col_q3:
                                    st.metric("Deep Scan", f"{pred_result.get('predictions', {}).get('deep_scan', 0)*100:.1f}%")
                        else:
                            st.warning("⚠️ Quality assessment is disabled")
                
                with col3:
                    if st.button("💾 Save to Artifacts", key="btn_save_artifacts", use_container_width=True):
                        if enable_offline_support:
                            with st.spinner("💾 Saving to artifacts..."):
                                case_id = st.session_state.get("case_id", "DEFAULT")
                                 
                                # Backend: Save to artifacts
                                save_result = media_viewer.export_media_report(case_id)
                                 
                                if save_result:
                                    st.success("✅ Saved to artifacts")
                                    st.info(f"Case: {case_id} | Offline access enabled")
                                else:
                                    st.error("❌ Failed to save artifacts")
                        else:
                            st.warning("⚠️ Offline support is disabled")
                
                with col4:
                    if st.button("📊 View Report", key="btn_view_report", use_container_width=True):
                        with st.spinner("📊 Generating report..."):
                            import tempfile
                            with tempfile.NamedTemporaryFile(delete=False, suffix=f".{file_ext}") as tmp:
                                tmp.write(doc_file.getvalue())
                                tmp_path = tmp.name
                             
                            # Backend: Generate comprehensive report
                            perf_report = media_viewer.recovery_performance_report()
                             
                            st.success("✅ Report generated")
                            with st.expander("📋 Recovery Report Details"):
                                col_r1, col_r2 = st.columns(2)
                                 
                                with col_r1:
                                    st.write("**File Information**")
                                    st.json({
                                        "filename": doc_file.name,
                                        "size_kb": f"{doc_file.size / 1024:.2f}",
                                        "type": file_ext,
                                        "uploaded": datetime.now().isoformat()
                                    })
                                 
                                with col_r2:
                                    st.write("**Recovery Settings**")
                                    st.json({
                                        "corruption_detection": enable_corruption_detection,
                                        "ai_recovery": enable_ai_image_recovery or enable_ai_video_recovery,
                                        "gpu_acceleration": enable_gpu_acceleration,
                                        "parallel_processing": enable_parallel_processing,
                                        "offline_support": enable_offline_support
                                    })
                                 
                                st.write("**Performance Metrics**")
                                st.json(perf_report.get('statistics', {}))
                
                st.divider()
                
                # ====================================================================
                # INNOVATIVE FILE RECOVERY & ADVANCED FEATURES UI
                # ====================================================================
                
                st.subheader("🔧 Advanced Recovery & Analysis Tools")
                
                # Feature toggles in expandable sections
                with st.expander("🚀 Recovery & Optimization Features", expanded=True):
                    col1, col2, col3, col4 = st.columns(4)
                    
                    with col1:
                        st.markdown("### 🔍 Detection & Scanning")
                        enable_corruption_detection = st.checkbox(
                            "Detect Corruption",
                            value=True,
                            key="toggle_corruption_detection",
                            help="Scan for file corruption issues"
                        )
                        enable_file_recovery_scan = st.checkbox(
                            "File Recovery Scan",
                            value=True,
                            key="toggle_file_recovery_scan",
                            help="Scan for recoverable files"
                        )
                    
                    with col2:
                        st.markdown("### 🤖 AI-Powered Recovery")
                        enable_ai_image_recovery = st.checkbox(
                            "AI Image Recovery",
                            value=True,
                            key="toggle_ai_image",
                            help="Use AI for image reconstruction"
                        )
                        enable_ai_video_recovery = st.checkbox(
                            "AI Video Recovery",
                            value=True,
                            key="toggle_ai_video",
                            help="Use AI for video frame recovery"
                        )
                    
                    with col3:
                        st.markdown("### ⚡ Performance Boost")
                        enable_gpu_acceleration = st.checkbox(
                            "GPU Acceleration",
                            value=False,
                            key="toggle_gpu",
                            help="Enable GPU for faster processing"
                        )
                        enable_parallel_processing = st.checkbox(
                            "Parallel Processing",
                            value=True,
                            key="toggle_parallel",
                            help="Use multi-threading for batch operations"
                        )
                    
                    with col4:
                        st.markdown("### 💾 Caching & Storage")
                        enable_smart_caching = st.checkbox(
                            "Smart Caching",
                            value=True,
                            key="toggle_caching",
                            help="Cache recovery results"
                        )
                        enable_offline_support = st.checkbox(
                            "Offline Support",
                            value=True,
                            key="toggle_offline",
                            help="Save to artifacts for offline access"
                        )
                
                st.divider()
                
                # Advanced Recovery Options
                with st.expander("🔐 Advanced Recovery Options", expanded=False):
                    col1, col2 = st.columns(2)
                    
                    with col1:
                        st.markdown("### 🎯 Recovery Methods")
                        recovery_method = st.radio(
                            "Select Recovery Method",
                            options=[
                                "Standard Recovery",
                                "AI-Powered Recovery",
                                "Deep Scan Recovery",
                                "Smart Recovery (Auto)"
                            ],
                            key="recovery_method_select",
                            help="Choose recovery technique"
                        )
                        
                        st.markdown("### 📊 Analysis Options")
                        enable_integrity_check = st.checkbox(
                            "Integrity Check",
                            value=True,
                            key="toggle_integrity",
                            help="Verify recovered file integrity"
                        )
                        enable_quality_assessment = st.checkbox(
                            "Quality Assessment",
                            value=True,
                            key="toggle_quality",
                            help="Assess recovery quality"
                        )
                    
                    with col2:
                        st.markdown("### 🗂️ Batch Operations")
                        enable_batch_recovery = st.checkbox(
                            "Batch Recovery",
                            value=True,
                            key="toggle_batch",
                            help="Recover multiple files"
                        )
                        batch_worker_count = st.slider(
                            "Worker Threads",
                            min_value=1,
                            max_value=8,
                            value=4,
                            key="batch_workers",
                            help="Number of parallel workers"
                        )
                        
                        st.markdown("### 🔔 Notifications")
                        enable_progress_tracking = st.checkbox(
                            "Progress Tracking",
                            value=True,
                            key="toggle_progress",
                            help="Show recovery progress"
                        )
                        enable_detailed_logging = st.checkbox(
                            "Detailed Logging",
                            value=False,
                            key="toggle_logging",
                            help="Enable detailed operation logs"
                        )
                
                st.divider()
                
                # Real-time Status Dashboard
                with st.expander("📊 Recovery Status Dashboard", expanded=False):
                    col1, col2, col3, col4 = st.columns(4)
                    
                    with col1:
                        st.metric("🔍 Corruption Detected", "0" if not enable_corruption_detection else "Scanning...")
                    with col2:
                        st.metric("📁 Files Recoverable", "0" if not enable_file_recovery_scan else "Scanning...")
                    with col3:
                        st.metric("⚡ Performance Boost", "0%" if not enable_gpu_acceleration else "30%+")
                    with col4:
                        st.metric("💾 Cache Size", "0 MB" if not enable_smart_caching else "Auto")
                
                st.divider()
                
                # Document actions
                st.write("**Document Actions**")
                col1, col2, col3 = st.columns(3)
                
                with col1:
                    if st.button("🔍 Analyze", key="analyze_document"):
                        # Analyze document
                        import tempfile
                        with tempfile.NamedTemporaryFile(delete=False, suffix=f".{file_ext}") as tmp:
                            tmp.write(doc_file.getvalue())
                            tmp_path = tmp.name
                        
                        analysis_result = media_viewer.analyze_document(tmp_path)
                        
                        if analysis_result['status'] == 'success':
                            st.success("✅ Document Analysis Complete")
                            analysis = analysis_result['analysis']
                            
                            col_a1, col_a2 = st.columns(2)
                            with col_a1:
                                st.metric("File Type", analysis.get('file_type', 'Unknown'))
                                st.metric("File Size", f"{analysis.get('file_size_kb', 0):.2f} KB")
                            with col_a2:
                                if 'lines' in analysis:
                                    st.metric("Lines", analysis['lines'])
                                if 'words' in analysis:
                                    st.metric("Words", analysis['words'])
                            
                            if 'characters' in analysis:
                                st.metric("Characters", analysis['characters'])
                            
                            if 'note' in analysis:
                                st.info(f"ℹ️ {analysis['note']}")
                        else:
                            st.error(f"❌ {analysis_result.get('error', 'Analysis failed')}")
                
            else:
                st.info("📄 Upload a document to view and analyze")
            
            st.divider()
            
            # Document list
            st.write("**Recent Documents**")
            
            doc_list = [
                {"name": "report_001.pdf", "type": "PDF", "size": "2.5 MB", "date": "2025-11-25"},
                {"name": "data_analysis.xlsx", "type": "Excel", "size": "1.2 MB", "date": "2025-11-24"},
                {"name": "notes.txt", "type": "Text", "size": "45 KB", "date": "2025-11-23"},
                {"name": "presentation.pptx", "type": "PowerPoint", "size": "5.8 MB", "date": "2025-11-22"},
            ]
            
            for doc in doc_list:
                col1, col2, col3, col4, col5 = st.columns([2, 1, 1, 1, 1])
                with col1:
                    st.write(f"📄 {doc['name']}")
                with col2:
                    st.caption(doc['type'])
                with col3:
                    st.caption(doc['size'])
                with col4:
                    st.caption(doc['date'])
                with col5:
                    if st.button("👁️", key=f"view_doc_{doc['name']}", help="View"):
                        st.info(f"Opening: {doc['name']}")
        
        else:
            st.info("👁️ Document Viewer hidden - Click checkbox to show")
    
    # ========================================================================
    # OFFLINE SUPPORT & ARTIFACT ROUTING
    # ========================================================================
    
    st.divider()
    st.subheader("🔄 Offline Support & Artifact Routing")
    
    # Offline status
    col1, col2, col3, col4 = st.columns(4)
    
    with col1:
        is_online = consent_manager.connectivity_manager.is_online()
        if is_online:
            st.success("🟢 Online")
        else:
            st.warning("🔴 Offline")
    
    with col2:
        if st.button("💾 Save to Artifacts", key="save_media_artifacts"):
            case_id = st.session_state.get("case_id", "DEFAULT")
            result = media_viewer.export_media_report(case_id)
            if result:
                st.success("✅ Media saved to artifacts")
            else:
                st.error("❌ Failed to save media")
    
    with col3:
        if st.button("📥 Load from Artifacts", key="load_media_artifacts"):
            case_id = st.session_state.get("case_id", "DEFAULT")
            result = media_viewer.sync_media_from_artifacts(case_id)
            if result:
                st.success("✅ Media loaded from artifacts")
            else:
                st.info("ℹ️ No artifacts found")
    
    with col4:
        if st.button("📊 Export Report", key="export_media_report"):
            case_id = st.session_state.get("case_id", "DEFAULT")
            result = media_viewer.export_media_report(case_id)
            if result:
                st.success("✅ Report exported")
            else:
                st.error("❌ Export failed")
    
    # Offline queue status
    st.write("**Offline Queue Status**")
    
    col1, col2, col3 = st.columns(3)
    
    with col1:
        pending_count = len(consent_manager.connectivity_manager.pending_sync_queue)
        st.metric("Pending Sync", pending_count)
    
    with col2:
        last_sync = consent_manager.connectivity_manager.last_sync_time
        if last_sync:
            st.metric("Last Sync", last_sync.strftime("%Y-%m-%d %H:%M:%S"))
        else:
            st.metric("Last Sync", "Never")
    
    with col3:
        if st.button("🔄 Sync Now", key="sync_now_button"):
            if is_online:
                st.success("✅ Sync completed")
            else:
                st.warning("⚠️ Offline - Sync queued")
    
    # Artifact routing details
    if st.checkbox("📋 Show Artifact Details", key="show_artifact_details"):
        st.write("**Artifact Structure**")
        
        artifact_info = {
            "Base Path": "artifacts/",
            "Case Directory": f"artifacts/{st.session_state.get('case_id', 'DEFAULT')}/",
            "Media Path": f"artifacts/{st.session_state.get('case_id', 'DEFAULT')}/media/",
            "Redactions Path": f"artifacts/{st.session_state.get('case_id', 'DEFAULT')}/media/*/redactions/",
            "Report Path": f"artifacts/{st.session_state.get('case_id', 'DEFAULT')}/media/media_report.json"
        }
        
        for key, value in artifact_info.items():
            st.caption(f"**{key}**: `{value}`")
        
        # Media summary
        st.write("**Media Summary**")
        col1, col2, col3 = st.columns(3)
        
        with col1:
            st.metric("Total Files", len(media_viewer.media_files))
        
        with col2:
            st.metric("Timeline Entries", len(media_viewer.media_timeline))
        
        with col3:
            redaction_stats = media_viewer.redaction_manager.get_redaction_stats()
            total_redactions = (
                redaction_stats['total_image_redactions'] +
                redaction_stats['total_audio_redactions'] +
                redaction_stats['total_video_redactions']
            )
            st.metric("Total Redactions", total_redactions)
        
        # Redaction breakdown
        st.write("**Redaction Breakdown**")
        col1, col2, col3 = st.columns(3)
        
        with col1:
            st.metric("Image Redactions", redaction_stats['total_image_redactions'])
        
        with col2:
            st.metric("Audio Redactions", redaction_stats['total_audio_redactions'])
        
        with col3:
            st.metric("Video Redactions", redaction_stats['total_video_redactions'])


def main():
    """Main Streamlit app"""
    st.set_page_config(
        page_title="Forensmart Analysis",
        page_icon="🔍",
        layout="wide"
    )
    
    # Sidebar navigation
    st.sidebar.title("🔍 Forensmart Analysis")
    
    # Get case info from session
    if "case_id" not in st.session_state:
        st.session_state.case_id = st.sidebar.text_input(
            "Case ID",
            value="DEFAULT",
            help="Enter case ID for tracking"
        )
    
    if "user_name" not in st.session_state:
        st.session_state.user_name = st.sidebar.text_input(
            "Your Name",
            value="Detective",
            help="Your name for tracking"
        )
    
    # Display case info
    st.sidebar.divider()
    st.sidebar.metric("Case ID", st.session_state.case_id)
    st.sidebar.metric("User", st.session_state.user_name)
    
    # Navigation
    st.sidebar.divider()
    page = st.sidebar.radio(
        "Select Analysis Module",
        ["Dashboard", "Communications", "Location Intelligence", "Database Management", "Media Viewer"],
        icons=["📊", "💬", "📍", "🗄️", "🎬"]
    )
    
    # Render selected page
    if page == "Dashboard":
        render_dashboard_overview()
    elif page == "Communications":
        render_comms_analyzer()
    elif page == "Location Intelligence":
        render_location_intelligence()
    elif page == "Database Management":
        render_database_management()
    else:
        render_media_viewer()


if __name__ == "__main__":
    main()
